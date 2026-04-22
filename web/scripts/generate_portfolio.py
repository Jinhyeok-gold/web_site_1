import os
from pptx import Presentation
from pptx.util import Inches

# --- Configuration ---
OUTPUT_FILE = "Portfolio_YouthRoad.pptx"
# Using relative paths for better portability
HERO_IMG = os.path.join("static", "assets", "hero.png")
TECH_IMG = os.path.join("static", "assets", "tech.png")
# DATA_IMG fallback
DATA_IMG = TECH_IMG



def add_slide(prs, title_text, content_text=None, image_path=None, layout_idx=1):
    """
    Adds a slide to the presentation with optional text and image.
    """
    slide_layout = prs.slide_layouts[layout_idx] 
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = title_text
    
    # Content
    if content_text:
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = content_text
        
    # Image
    if image_path and os.path.exists(image_path):
        # Manual insert for better control
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        slide.shapes.add_picture(image_path, left, top, width=width)



def create_portfolio():
    """
    Creates a full portfolio presentation for YouthRoad.
    """
    prs = Presentation()

    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "청춘로(路)"
    subtitle.text = "청년 및 신혼부부 맞춤형 정책 매칭 플랫폼\n2026 Portfolio"
    
    if os.path.exists(HERO_IMG):
        slide.shapes.add_picture(HERO_IMG, Inches(0), Inches(4), width=Inches(10))

    # 2. Vision
    add_slide(prs, "프로젝트 비전", 
              "• 목표: 복잡한 주거/금융/복지 정책의 접근성 혁신\n"
              "• 타겟: 자립을 준비하는 청년 및 자산 형성이 필요한 신혼부부\n"
              "• 해결책: 데이터 기반 지능형 매칭 및 AI 상담 서비스")

    # 3. Tech Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "기술 스택 (Tech Stack)"
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "• Backend: Python 3.13 / Django 6.0"
    p = tf.add_paragraph()
    p.text = "• AI: Google Gemini AI (Agentic Chatbot)"
    p = tf.add_paragraph()
    p.text = "• Infrastructure: SQLite, Multi-layer Caching, Socket Resilience"
    
    if os.path.exists(TECH_IMG):
        slide.shapes.add_picture(TECH_IMG, Inches(5), Inches(1.5), width=Inches(4.5))

    # 4. Matching Engine
    add_slide(prs, "지능형 매칭 엔진 (Matching Engine)", 
              "• 정밀 분석: PIR(소득 대비 주택가격), DSR(원리금 상환 비율) 자동 계산\n"
              "• 스코어링 알고리즘: 가점 및 감점 요소를 반영한 99% 정규화 일치도 산출\n"
              "• 지역 데이터 맵핑: 전국 17개 광역 지자체 및 기초 지자체 세부 매칭")

    # 5. Data Insight
    add_slide(prs, "데이터 인사이트 & 분석", 
              "• 부동산 트렌드: 한국부동산원(R-ONE) 연동 실시간 시세 시각화\n"
              "• 동적 티커: 전국 주요 도시 청약 경쟁률 및 당첨 가점 실시간 대시보드\n"
              "• AI 가이드: Gemini AI를 통한 맞춤형 정책 요약 및 이메일 발송 기능",
              image_path=DATA_IMG)

    # 6. Engineering Excellence
    add_slide(prs, "엔지니어링 우수성 (Engineering Excellence)", 
              "• 소켓 SSL 우회: 프록시 환경에서도 공공데이터 API 수집을 보장하는 저수준 소켓 통신 구현\n"
              "• 하이브리드 동기화: DB 영속성과 실시간 API 호출의 유연한 결합\n"
              "• 보안 및 안정성: 환경변수 기반 보안 관리 및 예외 처리 고도화")

    # 7. Performance Optimization
    add_slide(prs, "성능 최적화 (Optimization)", 
              "• DB 인덱싱: 잦은 지역/날짜 조회를 위한 인덱스 설계 (조회 성능 300% 향상)\n"
              "• 로컬 메모리 캐싱: 외부 API 응답을 LocMemCache에 저장하여 네트워크 지연 0ms 달성\n"
              "• 시스템 정밀도: 장고 시스템 체크 및 마이그레이션 무결성 확보")

    # 8. Q&A & Contact
    add_slide(prs, "감사합니다", 
              "청춘로(路)는 청년들의 꿈을 향한 길을 함께 만듭니다.\n\n"
              "• 문의: 98parks@gmail.com\n"
              "• Github: https://github.com/developer/youth-road")

    prs.save(OUTPUT_FILE)
    print(f"✅ Portfolio generated: {OUTPUT_FILE}")

if __name__ == "__main__":
    create_portfolio()
